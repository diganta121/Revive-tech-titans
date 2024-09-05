from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Automated Disaster Relief System with Payload Delivery"
subtitle.text = "An Advanced Approach to Delivering Supplies in Hazardous Areas\nPresented by: Your Name"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.shapes.placeholders[1].text_frame
title_2.text = "Introduction"
content_2.text = "Problem Statement: Delivering supplies to isolated or dangerous areas is a challenge."

# Add bullet points
bullet_2 = content_2.add_paragraph()
bullet_2.text = "Need for a robust, reliable system for payload delivery during disasters."

# Slide 3: System Overview
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.shapes.placeholders[1].text_frame
title_3.text = "System Overview"
content_3.text = "Key Components:"

# Add bullet points
content_3.add_paragraph("Payload Handling and Securing Mechanism")
content_3.add_paragraph("Navigation and Obstacle Detection")
content_3.add_paragraph("GPS-denied Environment Navigation")
content_3.add_paragraph("Power Management")

# Slide 4: Design for Adverse Conditions
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.shapes.placeholders[1].text_frame
title_4.text = "Design for Adverse Conditions"
content_4.text = "Challenges:"

# Add bullet points
content_4.add_paragraph("Harsh weather, difficult terrains")
content_4.add_paragraph("Securing and safely delivering payloads")

# Add sub-bullet points for solutions
content_4.add_paragraph("Solutions:").level = 0
content_4.add_paragraph("Weatherproof materials and structure").level = 1
content_4.add_paragraph("Advanced payload securing and release mechanisms").level = 1

# Slide 5: Power Management and Motor Control
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.shapes.placeholders[1].text_frame
title_5.text = "Power Management and Motor Control"
content_5.text = "Power Considerations:"

# Add bullet points
content_5.add_paragraph("Efficient use of battery or hybrid systems")
content_5.add_paragraph("Solar panels for extended operations")

# Add sub-bullet points for motor control
content_5.add_paragraph("Motor Control:").level = 0
content_5.add_paragraph("Precision control for navigation in tough terrains").level = 1

# Slide 6: Obstacle Detection and Navigation
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.shapes.placeholders[1].text_frame
title_6.text = "Obstacle Detection and Navigation"
content_6.text = "Obstacle Detection Methods:"

# Add bullet points
content_6.add_paragraph("LIDAR, Ultrasonic, and Infrared sensors")
content_6.add_paragraph("Cameras for vision-based navigation")

# Add sub-bullet points for algorithms
content_6.add_paragraph("Algorithms:").level = 0
content_6.add_paragraph("Real-time obstacle avoidance").level = 1
content_6.add_paragraph("Path planning using A* or Dijkstra algorithms").level = 1

# Slide 7: GPS-Denied Navigation
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.shapes.placeholders[1].text_frame
title_7.text = "GPS-Denied Navigation"
content_7.text = "Challenges:"

# Add bullet points
content_7.add_paragraph("Operating in areas with no GPS signal")

# Add sub-bullet points for solutions
content_7.add_paragraph("Solutions:").level = 0
content_7.add_paragraph("SLAM (Simultaneous Localization and Mapping)").level = 1
content_7.add_paragraph("Use of visual markers, RFID, or IMUs for localization").level = 1

# Slide 8: Payload Delivery Mechanism
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.shapes.placeholders[1].text_frame
title_8.text = "Payload Delivery Mechanism"
content_8.text = "Accuracy in Delivery:"

# Add bullet points
content_8.add_paragraph("Air or ground-based payload drop system")
content_8.add_paragraph("Parachute or tethered delivery for airborne systems")

# Add sub-bullet points for payload types
content_8.add_paragraph("Payload Types:").level = 0
content_8.add_paragraph("Medical supplies, food, communication devices, etc.").level = 1

# Slide 9: Conclusion
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
content_9 = slide_9.shapes.placeholders[1].text_frame
title_9.text = "Conclusion"
content_9.text = "Summary:"

# Add bullet points
content_9.add_paragraph("System is capable of handling challenging environments")
content_9.add_paragraph("Accurate and safe delivery using cutting-edge technology")
content_9.add_paragraph("Reliable obstacle detection and GPS-denied navigation")

# Save the presentation
prs.save("Automated_Disaster_Relief_System.pptx")

print("Presentation created successfully.")

